#!/usr/bin/env python3
"""pbox-deck -- build a .pptx from a simple JSON spec, using python-pptx.

Usage:
    pbox-deck --spec deck.json --out deck.pptx [--brand default]

Spec (JSON):
{
  "title": "Quarterly review",
  "subtitle": "Q2 2026",
  "author": "Your Company",
  "slides": [
    {"title": "Highlights", "bullets": ["Revenue up 12%", "Two new clients"], "notes": "speaker notes"},
    {"title": "Plan", "text": "A short paragraph of body text.", "notes": "..."}
  ]
}

Local, offline, no API. _DECK_BUILDER_V1
"""
import argparse
import json
import os
import sys


def _load_brand(brand, here):
    path = os.path.join(here, "brand-profiles.json")
    try:
        with open(path) as f:
            profiles = json.load(f)
        return profiles.get(brand) or profiles.get("default") or {}
    except Exception:
        return {}


def build(spec, out_path, brand_name="default"):
    try:
        from pptx import Presentation
        from pptx.util import Pt, Inches
        from pptx.dml.color import RGBColor
    except ImportError:
        sys.stderr.write("python-pptx is not installed. Run: pip3 install python-pptx\n")
        sys.exit(2)

    here = os.path.dirname(os.path.abspath(__file__))
    brand = _load_brand(brand_name, here)
    accent = brand.get("accent_rgb")  # e.g. [123, 104, 238]

    prs = Presentation()

    # Title slide (layout 0)
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = str(spec.get("title", "Untitled"))
    if s.placeholders and len(s.placeholders) > 1:
        sub = spec.get("subtitle") or spec.get("author") or ""
        s.placeholders[1].text = str(sub)
    if accent:
        try:
            s.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(*accent)
        except Exception:
            pass

    # Content slides (layout 1: title + content)
    for slide in spec.get("slides", []):
        cs = prs.slides.add_slide(prs.slide_layouts[1])
        cs.shapes.title.text = str(slide.get("title", ""))
        body = None
        for ph in cs.placeholders:
            if ph.placeholder_format.idx == 1:
                body = ph
                break
        if body is not None:
            tf = body.text_frame
            tf.clear()
            bullets = slide.get("bullets")
            if bullets:
                for i, b in enumerate(bullets):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = str(b)
                    p.level = 0
            elif slide.get("text"):
                tf.paragraphs[0].text = str(slide["text"])
        if accent:
            try:
                cs.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(*accent)
            except Exception:
                pass
        # speaker notes
        notes = slide.get("notes")
        if notes:
            cs.notes_slide.notes_text_frame.text = str(notes)

    prs.save(out_path)
    return len(spec.get("slides", [])) + 1


def main():
    ap = argparse.ArgumentParser(prog="pbox-deck", description="Build a .pptx from a JSON spec.")
    ap.add_argument("--spec", required=True, help="path to the deck JSON spec")
    ap.add_argument("--out", required=True, help="output .pptx path")
    ap.add_argument("--brand", default="default", help="brand profile name (brand-profiles.json)")
    args = ap.parse_args()

    with open(args.spec) as f:
        spec = json.load(f)
    n = build(spec, args.out, args.brand)
    print(f"[deck-builder] wrote {args.out} ({n} slides)")


if __name__ == "__main__":
    main()
